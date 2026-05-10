"""
Generate Z1 2-min share PPTX matching the Beamer PDF layout exactly.
5 slides: Title / Architecture / Strategy / Results / Sim2Sim
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ──
BASE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(BASE, "sources")

# ── Colors (matching Beamer theme) ──
DKBLUE = RGBColor(0, 60, 113)
NVGREEN = RGBColor(118, 185, 0)
WHITE = RGBColor(255, 255, 255)
LGRAY = RGBColor(240, 240, 240)
GRAY = RGBColor(120, 120, 120)
RED = RGBColor(200, 50, 50)
BLACK = RGBColor(30, 30, 30)

# ── Slide dimensions (16:9) ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=12,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_title_bar(slide, title_text):
    """Top dark blue bar with white title"""
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), DKBLUE)
    add_textbox(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.6),
                title_text, font_size=22, color=WHITE, bold=True)
    # Green accent line
    add_rect(slide, 0, Inches(0.85), SLIDE_W, Inches(0.04), NVGREEN)
    return bar


def add_block(slide, left, top, width, height, title, items, title_color=WHITE,
              bg_color=LGRAY, title_bg=NVGREEN, font_size=11):
    """Beamer-style block with green title + gray body"""
    # Title bar
    t_shape = add_rounded_rect(slide, left, top, width, Inches(0.35), title_bg)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.02), width - Inches(0.3), Inches(0.3),
                title, font_size=font_size, color=title_color, bold=True)
    # Body
    body_top = top + Inches(0.38)
    body_h = height - Inches(0.38)
    b_shape = add_rounded_rect(slide, left, body_top, width, body_h, bg_color)
    # Items text
    txBox = slide.shapes.add_textbox(left + Inches(0.15), body_top + Inches(0.08),
                                     width - Inches(0.3), body_h - Inches(0.15))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size - 1)
        p.font.color.rgb = BLACK
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(2)
    return b_shape


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)
    else:
        add_textbox(slide, left, top, Inches(2), Inches(0.3),
                    f"[missing: {os.path.basename(path)}]", font_size=9, color=RED)


# ================================================================
# Slide 1: Title
# ================================================================
slide = add_blank_slide(prs)
# Full dark blue background
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DKBLUE)
# Green accent line
add_rect(slide, 0, Inches(3.2), SLIDE_W, Inches(0.05), NVGREEN)

# Left: title text
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(6.5), Inches(0.8),
            "MagicBot Z1 12DOF", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(2.3), Inches(6.5), Inches(0.6),
            "基于强化学习的双足机器人运动控制", font_size=22, color=LGRAY)
add_textbox(slide, Inches(0.8), Inches(3.5), Inches(6.5), Inches(0.4),
            "Isaac Lab + PPO  |  16384 envs  |  Curriculum Learning", font_size=13, color=GRAY)
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(0.3),
            "MagicLab  |  2026.05", font_size=12, color=GRAY)

# Right: robot image
add_image_safe(slide, os.path.join(SRC, "z1_mujoco.png"),
               Inches(8.0), Inches(1.2), width=Inches(4.5))

# ================================================================
# Slide 2: System Architecture
# ================================================================
slide = add_blank_slide(prs)
add_title_bar(slide, "系统架构 — Two-Platform Training Pipeline")

# Left column: two blocks
col1_left = Inches(0.4)
col1_w = Inches(6.2)

add_block(slide, col1_left, Inches(1.2), col1_w, Inches(2.3),
          "RTX 6000 远程训练",
          ["● 8x RTX 6000 GPU (85GB VRAM)",
           "● 16384 并行仿真环境",
           "● Isaac Lab + rsl_rl PPO",
           "● 自动化 5-Phase Pipeline"])

add_block(slide, col1_left, Inches(3.8), col1_w, Inches(1.9),
          "本地分析 & 验证",
          ["● MuJoCo Sim2Sim 验证",
           "● 学习曲线 & 过拟合分析",
           "● 键盘实时操控测试"])

# Right column: 5-Phase flow
col2_left = Inches(7.0)
col2_w = Inches(5.8)

phases = ["P1 平地粗训", "P2 平地精调", "P3 缓坡地形", "P3b 中等地形", "P4 复杂地形"]
phase_top = Inches(1.2)
box_h = Inches(0.55)
gap = Inches(0.15)

# Block background
add_block(slide, col2_left, Inches(1.2), col2_w, Inches(4.5),
          "5-Phase Curriculum", [], title_bg=NVGREEN)

for i, phase in enumerate(phases):
    y = Inches(1.75) + (box_h + gap) * i
    box = add_rounded_rect(slide, col2_left + Inches(1.2), y, Inches(3.4), box_h,
                           RGBColor(235, 250, 220))
    box.line.color.rgb = NVGREEN
    box.line.width = Pt(1)
    add_textbox(slide, col2_left + Inches(1.5), y + Inches(0.08), Inches(3), Inches(0.4),
                phase, font_size=13, color=DKBLUE, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < len(phases) - 1:
        add_textbox(slide, col2_left + Inches(2.7), y + box_h - Inches(0.05),
                    Inches(0.5), Inches(0.2), "▼", font_size=10, color=DKBLUE,
                    alignment=PP_ALIGN.CENTER)

# ================================================================
# Slide 3: Training Strategy
# ================================================================
slide = add_blank_slide(prs)
add_title_bar(slide, "训练策略 — Reward 设计 & PPO")

# Left: reward items
left_col = Inches(0.4)
left_w = Inches(4.5)

# Reward table - positive
add_textbox(slide, left_col, Inches(1.1), left_w, Inches(0.3),
            "激励项", font_size=13, color=DKBLUE, bold=True)

pos_rewards = [
    ("XY 速度跟踪", "w = 1.0"),
    ("角速度跟踪", "w = 0.5"),
    ("存活奖励", "w = 0.15"),
    ("足部接触时序", "w = 0.5"),
    ("足部摆动高度", "w = 1.0"),
]
y = Inches(1.45)
for name, weight in pos_rewards:
    add_textbox(slide, left_col + Inches(0.15), y, Inches(2.5), Inches(0.22),
                f"● {name}", font_size=10, color=BLACK)
    add_textbox(slide, left_col + Inches(2.8), y, Inches(1.5), Inches(0.22),
                weight, font_size=10, color=DKBLUE, alignment=PP_ALIGN.RIGHT)
    y += Inches(0.25)

# Reward table - negative
y += Inches(0.15)
add_textbox(slide, left_col, y, left_w, Inches(0.3),
            "惩罚项", font_size=13, color=RED, bold=True)
y += Inches(0.35)

neg_rewards = [
    ("Z 轴速度", "w = -2.0"),
    ("身体姿态偏移", "w = -5.0"),
    ("基座高度偏差", "w = -10.0"),
    ("能量消耗", "w = -2e-5"),
    ("足部滑动", "w = -0.2"),
    ("动作变化率", "w = -0.1"),
]
for name, weight in neg_rewards:
    add_textbox(slide, left_col + Inches(0.15), y, Inches(2.5), Inches(0.22),
                f"● {name}", font_size=10, color=BLACK)
    add_textbox(slide, left_col + Inches(2.8), y, Inches(1.5), Inches(0.22),
                weight, font_size=10, color=RED, alignment=PP_ALIGN.RIGHT)
    y += Inches(0.25)

# Right: reward decomposition image + PPO config
right_col = Inches(5.2)
right_w = Inches(7.6)

add_image_safe(slide, os.path.join(SRC, "reward_decomposition.png"),
               right_col, Inches(1.1), width=Inches(7.6), height=Inches(3.5))

add_textbox(slide, right_col, Inches(4.7), right_w, Inches(0.25),
            "P2 Fine — 各 Reward Component 贡献趋势", font_size=9, color=GRAY,
            alignment=PP_ALIGN.CENTER)

# PPO config block
add_block(slide, right_col, Inches(5.1), right_w, Inches(1.8),
          "PPO 配置",
          ["网络: MLP (32, 32), Actor-Critic",
           "Learning Rate: 3e-4",
           "Entropy Coeff: 0.01",
           "GAE λ: 0.95"],
          font_size=10)

# ================================================================
# Slide 4: Training Results
# ================================================================
slide = add_blank_slide(prs)
add_title_bar(slide, "训练成果")

# Left: curriculum plot + demo GIF
left_col = Inches(0.4)
add_image_safe(slide, os.path.join(SRC, "curriculum_reward_trends.png"),
               left_col, Inches(1.1), width=Inches(6.5), height=Inches(3.2))
add_textbox(slide, left_col, Inches(4.35), Inches(6.5), Inches(0.25),
            "Curriculum — 多 Phase Reward 趋势", font_size=9, color=GRAY,
            alignment=PP_ALIGN.CENTER)

# Demo GIF
gif_path = os.path.join(SRC, "pipeline_demo.gif")
add_image_safe(slide, gif_path,
               Inches(2.0), Inches(4.7), width=Inches(3.2))
add_textbox(slide, Inches(2.0), Inches(7.0), Inches(3.2), Inches(0.25),
            "P1→P2 训练演示", font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

# Right: metrics + training scale
right_col = Inches(7.3)
right_w = Inches(5.5)

# Key metrics table
add_block(slide, right_col, Inches(1.1), right_w, Inches(2.8),
          "关键指标", [], title_bg=NVGREEN, font_size=10)

metrics = [
    ("P2 最佳 Reward", "49.68"),
    ("步态距离 (Sim2Sim)", "4.0m / 10s"),
    ("Sim2Sim 摔倒率", "0%"),
    ("本地测试距离", "12.0m / 25.5s"),
    ("本地摔倒率", "0.06%"),
]
y = Inches(1.65)
# Header
add_textbox(slide, right_col + Inches(0.2), y, Inches(2.8), Inches(0.25),
            "指标", font_size=10, color=DKBLUE, bold=True)
add_textbox(slide, right_col + Inches(3.2), y, Inches(2), Inches(0.25),
            "数值", font_size=10, color=DKBLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, right_col + Inches(0.1), y + Inches(0.28), Inches(5.2), Inches(0.01), NVGREEN)
y += Inches(0.35)

for name, val in metrics:
    add_textbox(slide, right_col + Inches(0.2), y, Inches(2.8), Inches(0.22),
                name, font_size=10, color=BLACK)
    add_textbox(slide, right_col + Inches(3.2), y, Inches(2), Inches(0.22),
                val, font_size=10, color=DKBLUE, bold=True, alignment=PP_ALIGN.CENTER)
    y += Inches(0.32)

# Training scale block
add_block(slide, right_col, Inches(4.3), right_w, Inches(1.8),
          "训练规模",
          ["● 16384 并行环境",
           "● 5-Phase 自动化 Pipeline",
           "● 训练 → 导出 → Sim2Sim 验证"],
          font_size=10)

# ================================================================
# Slide 5: Sim2Sim & Next Steps
# ================================================================
slide = add_blank_slide(prs)
add_title_bar(slide, "Sim2Sim 验证 & 下一步")

# Left: sim2sim table + GIF
left_col = Inches(0.4)
left_w = Inches(6.2)

add_block(slide, left_col, Inches(1.1), left_w, Inches(2.2),
          "Sim2Sim 验证结果", [], title_bg=NVGREEN, font_size=10)

sim2sim = [
    ("P1 Fine", "平地", "OK", NVGREEN),
    ("P2 Fine", "平地", "OK", NVGREEN),
    ("P3 Fine", "缓坡", "20 falls", RED),
    ("P3b Fine", "中等", "冻结", RED),
]
y = Inches(1.65)
# Header
for col, (text, x_off, w) in enumerate([
    ("Phase", Inches(0.2), Inches(1.5)),
    ("地形", Inches(1.8), Inches(1.5)),
    ("MuJoCo", Inches(3.5), Inches(2.2)),
]):
    add_textbox(slide, left_col + x_off, y, w, Inches(0.25),
                text, font_size=10, color=DKBLUE, bold=True,
                alignment=PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT)
add_rect(slide, left_col + Inches(0.1), y + Inches(0.28), Inches(5.8), Inches(0.01), NVGREEN)
y += Inches(0.35)

for phase, terrain, result, color in sim2sim:
    add_textbox(slide, left_col + Inches(0.2), y, Inches(1.5), Inches(0.22),
                phase, font_size=10, color=BLACK)
    add_textbox(slide, left_col + Inches(1.8), y, Inches(1.5), Inches(0.22),
                terrain, font_size=10, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left_col + Inches(3.5), y, Inches(2.2), Inches(0.22),
                result, font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    y += Inches(0.3)

# Sim2Sim broken GIF
gif_path2 = os.path.join(SRC, "sim2sim_broken.gif")
add_image_safe(slide, gif_path2,
               Inches(1.5), Inches(3.6), width=Inches(3.5))
add_textbox(slide, Inches(1.5), Inches(6.8), Inches(3.5), Inches(0.25),
            "P3 Sim2Sim 失败 (自动播放)", font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

# Right: findings + next steps
right_col = Inches(7.0)
right_w = Inches(5.8)

add_block(slide, right_col, Inches(1.1), right_w, Inches(1.8),
          "核心发现",
          ["● 平地策略 → MuJoCo 迁移成功",
           "● 地形策略 → Sim2Sim Gap",
           "● 原因：物理引擎差异（接触力/摩擦）"],
          font_size=10)

# Alert block (orange/red)
t_shape = add_rounded_rect(slide, right_col, Inches(3.2), right_w, Inches(0.35),
                           RGBColor(200, 80, 40))
add_textbox(slide, right_col + Inches(0.15), Inches(3.22), right_w - Inches(0.3), Inches(0.3),
            "下一步计划", font_size=11, color=WHITE, bold=True)
body_shape = add_rounded_rect(slide, right_col, Inches(3.58), right_w, Inches(3.0),
                              RGBColor(255, 240, 230))

steps = [
    ("1. 解决 Sim2Sim Gap", "观测空间对齐 + 物理参数标定"),
    ("2. Sim2Real 真机部署", "Domain Randomization + 降阶"),
    ("3. 复杂地形自适应行走", "越障、楼梯、不平地面"),
]
y = Inches(3.7)
for title, desc in steps:
    add_textbox(slide, right_col + Inches(0.2), y, right_w - Inches(0.4), Inches(0.25),
                title, font_size=11, color=DKBLUE, bold=True)
    add_textbox(slide, right_col + Inches(0.4), y + Inches(0.25), right_w - Inches(0.6), Inches(0.2),
                desc, font_size=9, color=GRAY)
    y += Inches(0.6)

# Footer
add_textbox(slide, Inches(0.4), Inches(7.1), Inches(5), Inches(0.25),
            "MagicBot Z1  |  MagicLab  |  2026.05", font_size=8, color=GRAY)

# ── Save ──
output = os.path.join(BASE, "main_v2.pptx")
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
